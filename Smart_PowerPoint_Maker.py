import io
import os
import re
import time
import requests
import tkinter as tk
from tkinter import ttk, messagebox
from tkinter.scrolledtext import ScrolledText
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt
import webbrowser

# برای فرمت‌دهی به لینک‌ها
from pptx.dml.color import RGBColor

# کتابخانه چت بات (GenAI)
from google import genai

# ---------------- کلاس Tooltip جهت راهنمایی سریع ----------------
class CreateToolTip:
    def __init__(self, widget, text=''):
        self.waittime = 500  # زمان تاخیر قبل نمایش Tooltip (میلی‌ثانیه)
        self.wraplength = 180  # حداکثر عرض Tooltip
        self.widget = widget
        self.text = text
        self.widget.bind("<Enter>", self.enter)
        self.widget.bind("<Leave>", self.leave)
        self.id = None
        self.tw = None

    def enter(self, event=None):
        self.schedule()

    def leave(self, event=None):
        self.unschedule()
        self.hidetip()

    def schedule(self):
        self.unschedule()
        self.id = self.widget.after(self.waittime, self.showtip)

    def unschedule(self):
        id_ = self.id
        self.id = None
        if id_:
            self.widget.after_cancel(id_)

    def showtip(self, event=None):
        x, y, cx, cy = self.widget.bbox("insert")
        x += self.widget.winfo_rootx() + 25
        y += self.widget.winfo_rooty() + 20
        self.tw = tk.Toplevel(self.widget)
        self.tw.wm_overrideredirect(True)
        self.tw.wm_geometry("+%d+%d" % (x, y))
        label = tk.Label(self.tw, text=self.text, justify='left',
                         background="#ffffe0", relief='solid', borderwidth=1,
                         wraplength=self.wraplength, font=("Tahoma", 10))
        label.pack(ipadx=1)

    def hidetip(self):
        if self.tw:
            self.tw.destroy()
            self.tw = None

# ---------------- توابع بررسی وضعیت و پردازش ----------------
def fetch_data(url, query):
    payload = {'query': query}
    try:
        response = requests.post(url, data=payload)
        response.raise_for_status()
        response_data = response.json()
        if response_data.get('error'):
            raise ValueError(f"Server error: {response_data['error']}")
        return response_data['data']
    except (requests.RequestException, ValueError) as e:
        print(f"Error fetching data: {e}")
        return None

def check_app_state():
    url = 'https://rash32.ir/python/micropython/mysql_proxy_pptx_python_app.php'
    query = "SELECT * FROM pptx_python_app"
    data = fetch_data(url, query)
    if data:
        filtered_data = [item for item in data if item.get('state') == '1']
        if filtered_data:
            return True
    return False

def is_rtl(text):
    return bool(re.search(r'[\u0600-\u06FF]', text))

def process_text(text):
    return text

def remove_bullets(text_frame):
    for paragraph in text_frame.paragraphs:
        pPr = paragraph._p.get_or_add_pPr()
        buNone = OxmlElement('a:buNone')
        pPr.append(buNone)

def advanced_process_content(text):
    lines = text.splitlines()
    processed_paragraphs = []
    for line in lines:
        clean_line = re.sub(r'\s+', ' ', line).strip()
        if not clean_line:
            continue
        bullet = False
        if clean_line.startswith("•"):
            bullet = True
            clean_line = clean_line.lstrip("•").strip()
        processed_paragraphs.append({"text": clean_line, "bullet": bullet})
    return processed_paragraphs

def add_runs_with_links(paragraph, text, font_size=20):
    url_regex = r'(https?://\S+)'
    parts = re.split(url_regex, text)
    trailing_chars = '.,؛:!?'
    for part in parts:
        if re.match(url_regex, part):
            actual_url = part
            trailing = ''
            while actual_url and actual_url[-1] in trailing_chars:
                trailing = actual_url[-1] + trailing
                actual_url = actual_url[:-1]
            run = paragraph.add_run()
            run.text = actual_url
            run.font.size = Pt(font_size)
            run.font.underline = True
            run.font.color.rgb = RGBColor(0, 0, 255)
            if trailing:
                run2 = paragraph.add_run()
                run2.text = trailing
                run2.font.size = Pt(font_size)
        else:
            run = paragraph.add_run()
            run.text = part
            run.font.size = Pt(font_size)

def parse_slides(text):
    blocks = re.split(r"(?=اسلاید\s*(?:شماره:)?\s*\d+)", text)
    slides = []
    for block in blocks:
        block = block.strip()
        if not block:
            continue
        m_title = re.search(r"عنوان\s+اسلاید\s*:\s*(.+)", block)
        title = m_title.group(1).strip() if m_title else ""
        m_content = re.search(r"محتوا\s*:\s*(.+)", block, re.DOTALL)
        content = m_content.group(1).strip() if m_content else ""
        slides.append({"title": title, "content": content})
    return slides

def create_pptx_in_memory(slides_data):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    
    for slide in slides_data:
        current_slide = prs.slides.add_slide(blank_slide_layout)
        
        # جعبه عنوان
        left_title = Inches(1)
        top_title = Inches(0.5)
        width_title = Inches(8)
        height_title = Inches(1)
        title_box = current_slide.shapes.add_textbox(left_title, top_title, width_title, height_title)
        tf_title = title_box.text_frame
        tf_title.text = ""
        p_title = tf_title.add_paragraph()
        p_title.text = process_text(slide["title"])
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.alignment = PP_ALIGN.CENTER
        remove_bullets(tf_title)
        tf_title.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        # جعبه محتوا
        left_content = Inches(1)
        top_content = Inches(1.7)
        width_content = Inches(8)
        height_content = Inches(4)
        content_box = current_slide.shapes.add_textbox(left_content, top_content, width_content, height_content)
        tf_content = content_box.text_frame
        tf_content.text = ""
        tf_content.word_wrap = True
        tf_content.margin_left = Inches(0.1)
        tf_content.margin_right = Inches(0.1)
        tf_content.margin_top = Inches(0.1)
        tf_content.margin_bottom = Inches(0.1)
        tf_content.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        paragraphs = advanced_process_content(process_text(slide["content"]))
        for para in paragraphs:
            p = tf_content.add_paragraph()
            p.space_after = Pt(6)
            p.alignment = PP_ALIGN.RIGHT if is_rtl(para["text"]) else PP_ALIGN.LEFT
            if para["bullet"]:
                p.level = 1
            add_runs_with_links(p, para["text"], font_size=20)
    
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

def upload_file(file_bytes, file_name, upload_url):
    file_content = file_bytes.getvalue()
    
    def post_request(url):
        file_obj = io.BytesIO(file_content)
        files = {
            "file": (file_name, file_obj, "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        }
        data = {"destination_folder": "store_html_report"}
        headers = {"User-Agent": "Mozilla/5.0", "Expect": ""}
        return requests.post(url, files=files, data=data, headers=headers, allow_redirects=False)
    
    response = post_request(upload_url)
    if response.status_code in (301, 302, 303, 307, 308):
        redirect_url = response.headers.get("Location")
        if redirect_url:
            print("Redirected to:", redirect_url)
            response = post_request(redirect_url)
    
    if response.status_code == 200:
        try:
            resp_data = response.json()
            download_link = resp_data.get("download_link", "").strip()
            print("Download link received:", download_link)
            return download_link
        except Exception as e:
            print("Error processing JSON response:", e)
            return response.text.strip()
    else:
        print("File upload failed. Status code:", response.status_code)
        print("Response text:", response.text)
        return None

def save_download_link(download_link, txt_file):
    with open(txt_file, 'w', encoding='utf-8') as f:
        f.write(download_link)
    print(f"Download link saved in file '{txt_file}'.")

text_tab_shown = False

# ---------------- واسط کاربری (GUI) با رعایت دقیق اصول UI/UX ----------------
def main():
    root = tk.Tk()
    root.title("پاورپوینت ساز هوشمند")
    root.geometry("1200x850")
    root.minsize(1100, 800)
    root.configure(bg="#1e1e2f")  # پس‌زمینه تاریک مدرن
    
    # ---------------- ایجاد هدر جذاب ----------------
    header_frame = tk.Frame(root, bg="#283655", height=60)
    header_frame.pack(side="top", fill="x")
    header_label = tk.Label(header_frame, text="پاورپوینت ساز هوشمند", bg="#283655", fg="white",
                            font=("B Nazanin", 24, "bold"))
    header_label.pack(pady=10)
    
    # نوار منو
    menu_bar = tk.Menu(root, bg="#283655", fg="white")
    file_menu = tk.Menu(menu_bar, tearoff=0, background="#3b4b69", foreground="white")
    file_menu.add_command(label="خروج", command=root.quit)
    menu_bar.add_cascade(label="فایل", menu=file_menu)
    help_menu = tk.Menu(menu_bar, tearoff=0, background="#3b4b69", foreground="white")
    help_menu.add_command(label="ورژن", command=lambda: messagebox.showinfo("درباره", 
        "پاورپوینت ساز هوشمند نسخه 1.0\nتوسعه توسط Ali Rashidi"))
    menu_bar.add_cascade(label="ورژن", menu=help_menu)
    root.config(menu=menu_bar)
    
    # استایل دهی
    style = ttk.Style()
    style.theme_use('clam')
    style.configure("TLabel", background="#1e1e2f", foreground="white", font=("Segoe UI", 12))
    style.configure("TButton", font=("Segoe UI", 12, "bold"), background="#283655", foreground="white")
    style.map("TButton", background=[('active', '#3b4b69')])
    style.configure("TFrame", background="#1e1e2f")
    
    # ---------------- ایجاد نوت‌بوک تب‌ها ----------------
    notebook = ttk.Notebook(root)
    notebook.pack(fill='both', expand=True, padx=10, pady=10)
    
    def add_context_menu(widget):
        context_menu = tk.Menu(widget, tearoff=0)
        context_menu.add_command(label="کپی", command=lambda: widget.event_generate("<<Copy>>"))
        context_menu.add_command(label="چسباندن", command=lambda: widget.event_generate("<<Paste>>"))
        context_menu.add_command(label="انتخاب همه", command=lambda: widget.event_generate("<<SelectAll>>"))
        def show_context(event):
            context_menu.tk_popup(event.x_root, event.y_root)
        widget.bind("<Button-3>", show_context)
        widget.bind("<Control-a>", lambda e: widget.event_generate("<<SelectAll>>"))
    
    # ---------------- تب ورودی متن ----------------
    manual_frame = ttk.Frame(notebook, padding=10)
    notebook.add(manual_frame, text="ورودی متن")
    
    # فیلد موضوع
    subject_frame = ttk.Frame(manual_frame)
    subject_frame.grid(row=0, column=0, sticky="ew", pady=(0,5))
    subject_label = ttk.Label(subject_frame, text="موضوع:")
    subject_label.pack(side="left", padx=(0,5))
    subject_entry = ttk.Entry(subject_frame, font=("Segoe UI", 11))
    subject_entry.pack(side="left", fill="x", expand=True)
    
    # فیلد تعداد اسلاید
    slide_count_frame = ttk.Frame(manual_frame)
    slide_count_frame.grid(row=1, column=0, sticky="ew", pady=(0,5))
    slide_count_label = ttk.Label(slide_count_frame, text="تعداد اسلاید:")
    slide_count_label.pack(side="left", padx=(0,5))
    slide_count_entry = ttk.Entry(slide_count_frame, font=("Segoe UI", 11))
    slide_count_entry.pack(side="left", fill="x", expand=True)
    
    # فیلد متن پرامت
    manual_label = ttk.Label(manual_frame, text="اول موضوع و و تعداد اسلاید را وارد کن و در ادامه دکمه (انتقال متن به چت بات) را بزن :", font=("Segoe UI", 12))
    manual_label.grid(row=2, column=0, sticky="w", pady=(0,5))
    
    text_input = ScrolledText(manual_frame, height=20, width=100, font=("Segoe UI", 11), wrap="word")
    text_input.grid(row=3, column=0, padx=10, pady=10, sticky="nsew")
    add_context_menu(text_input)
    manual_frame.grid_rowconfigure(3, weight=1)
    manual_frame.grid_columnconfigure(0, weight=1)
    
    original_sample_text = (
        "به عنوان یک متخصص در زمینه تولید و طراحی ارائه‌های پاورپوینت، با 40 سال تجربه واقعی و مثل مادری دلسوز، لطفاً یک پاورپوینت جامع، کاربردی، جذاب و حرفه‌ای و مرتبط با موضوع **[موضوع]** تهیه کن و دارای تعداد **[تعداد]** اسلاید باشد. \n"
        "در تهیه این ارائه، از تکنیک‌های پیشرفته مانند **Chain-of-Thought** برای تحلیل گام‌به‌گام و مستندسازی دقیق هر اسلاید استفاده کن تا روند تفکر و تصمیم‌گیری به صورت شفاف و مستند ارائه شود و خروجی اسمی از این تکنیک‌ها نباشد.\n\n\n\n"
        " *نکات کلیدی :* \n\n"
        "1️⃣ **استفاده از چندنمونه‌ای (Few-shot Learning):** نمونه‌های موفق را به عنوان مرجع در نظر بگیر.  \n"
        "2️⃣ **تقسیم‌بندی گام‌به‌گام:** مسئله را به بخش‌های کوچک تقسیم کن و در هر بخش نقش‌های مشخصی را تخصیص بده.  \n"
        "3️⃣ **انسجام و ارتباط اسلایدها:** تمامی اسلایدها باید از نظر محتوایی و طراحی یکپارچه باشند.  \n"
        "4️⃣ **استفاده از منابع معتبر:** در صورت نیاز، داده‌ها و اطلاعات را از منابع علمی استخراج کن.  \n"
        "5️⃣ **سبک و قالب حرفه‌ای:** از زبان ساده، دقیق و حرفه‌ای استفاده کن که برای دانش‌آموزان و علاقه‌مندان قابل فهم باشد.  \n"
        "6️⃣ **عدم استفاده از bullet points یا جداکننده‌های اضافی.**  \n"
        "7️⃣ **بازبینی و بهبود:** در صورت نیاز، بازخورد بده و کیفیت نهایی را بهبود بخش.  \n\n\n\n"
        " *ساختار خروجی پاورپوینت :*  \n"
        "---------------------------------------------------------\n"
        "اسلاید شماره: [شماره]\n"
        "عنوان اسلاید: [عنوان]\n"
        "محتوا: [متن اصلی اسلاید]\n"
        "---------------------------------------------------------\n\n\n\n"
        "📌 *توجه :*  \n"
        "✅ از معادل‌های فارسی یا ترجمه شده به انگلیسی استفاده کن چون خروجی بهم میریزد.\n"
        "✅ خروجی با زبان فارسی باشد و بدون هرگونه توضیح اضافه باشد.\n"
        "✅ از اضافه‌گویی پرهیز کن و محتوا بصورت تیتروار با کمی توضیحات باشد.\n"
        "✅ اطلاعات باید دقیق و علمی باشند.\n"
        "✅ ارائه باید یکپارچگی محتوایی و بصری داشته باشد.\n"
        "✅ تمامی توضیحات باید به زبان ساده، اما دقیق و علمی باشند.\n"
    )
    text_input.insert("1.0", original_sample_text)
    
    # فریم دکمه‌ها (برای دکمه‌های مربوط به تب ورودی متن)
    buttons_frame = ttk.Frame(manual_frame)
    buttons_frame.grid(row=4, column=0, pady=5, sticky="ew")
    
    clear_manual_button = ttk.Button(buttons_frame, text="پاکسازی متن ورودی", command=lambda: text_input.delete("1.0", "end"))
    clear_manual_button.pack(side='left', padx=5)
    CreateToolTip(clear_manual_button, "متن ورودی را پاک می‌کند")
    
    def copy_prompt_text():
        prompt = text_input.get("1.0", "end").strip()
        root.clipboard_clear()
        root.clipboard_append(prompt)
        messagebox.showinfo("کپی", "متن پرامت کپی شد.")
    
    copy_prompt_button = ttk.Button(buttons_frame, text="(برای ویرایش دستی آن)کپی متن پرامت", command=copy_prompt_text)
    copy_prompt_button.pack(side='left', padx=5)
    CreateToolTip(copy_prompt_button, "متن پرامت را کپی می‌کند")
    
    def paste_clipboard_text():
        try:
            clipboard_text = root.clipboard_get()
            text_input.delete("1.0", "end")
            text_input.insert("1.0", clipboard_text)
            messagebox.showinfo("انتقال", "متن از کلیپ بورد انتقال یافت.")
        except Exception as e:
            messagebox.showerror("خطا", "کلیپ بورد خالی یا در دسترس نیست.")
    
    paste_clipboard_button = ttk.Button(buttons_frame, text="انتقال متن از کلیپ بورد به ورودی متن", command=paste_clipboard_text)
    paste_clipboard_button.pack(side='left', padx=5)
    CreateToolTip(paste_clipboard_button, "متن کلیپ بورد را وارد می‌کند")
    
    # دکمه انتقال متن به پرامت چت بات با جایگزینی موضوع و تعداد
    def transfer_text_to_chatbot():
        subject = subject_entry.get().strip()
        slide_count = slide_count_entry.get().strip()
        manual_text = text_input.get("1.0", "end")
        updated_text = manual_text.replace("[موضوع]", subject).replace("[تعداد]", slide_count)
        chatbot_prompt_input.delete("1.0", "end")
        chatbot_prompt_input.insert("1.0", updated_text)
        messagebox.showinfo("انتقال متن", "متن پرامت با موضوع و تعداد به تب چت بات منتقل شد.")
    
    transfer_to_chatbot_button = ttk.Button(buttons_frame, text="انتقال متن ورودی به ورودی چت بات", command=transfer_text_to_chatbot)
    transfer_to_chatbot_button.pack(side='left', padx=5)
    CreateToolTip(transfer_to_chatbot_button, "متن پرامت (با موضوع و تعداد) را به تب چت بات منتقل می‌کند")
    
    # تعریف یک فریم جدید برای دکمه تولید و آپلود پاورپوینت
    manual_action_frame = ttk.Frame(manual_frame, padding=10)
    manual_action_frame.grid(row=5, column=0, sticky="ew")
    
    download_link_var = tk.StringVar()
    
    def process_input():
        root.config(cursor="wait")
        root.update()
        if not check_app_state():
            messagebox.showerror("غیر فعال", "برنامه فعلا در دسترس نیست.\nلطفاً بعداً امتحان کنید.")
            root.config(cursor="")
            return
        raw_text = text_input.get("1.0", "end").strip()
        if not raw_text:
            messagebox.showerror("خطا", "متن ورودی خالی است!")
            root.config(cursor="")
            return
        slides_data = parse_slides(raw_text)
        if not slides_data:
            messagebox.showerror("خطا", "فرمت ورودی نادرست است.\nلطفاً بررسی کنید.")
            root.config(cursor="")
            return
        pptx_bytes = create_pptx_in_memory(slides_data)
        timestamp = int(time.time())
        file_name = f"promptx_pitch_deck_{timestamp}.pptx"
        upload_url = "http://www.rash32.ir/python/micropython/file_uploader_proxy.php"
        download_link = upload_file(pptx_bytes, file_name, upload_url)
        root.config(cursor="")
        if download_link:
            download_link_var.set(download_link)
            messagebox.showinfo("موفقیت", f"پاورپوینت تولید و آپلود شد!\nلینک دانلود:\n{download_link}")
        else:
            messagebox.showerror("خطا", "آپلود فایل با مشکل مواجه شد.")
    
    generate_button = ttk.Button(manual_action_frame, text="تولید پاورپوینت و ساخت لینک دانلود آن", command=process_input)
    generate_button.pack(side='left', padx=5)
    CreateToolTip(generate_button, "پاورپوینت را تولید و آپلود می‌کند")
    
    download_link_label = ttk.Label(manual_action_frame, textvariable=download_link_var, foreground="cyan", cursor="hand2", font=("Segoe UI", 11, "underline"))
    download_link_label.pack(side='left', padx=10)
    download_link_label.bind("<Button-1>", lambda e: webbrowser.open(download_link_var.get()))
    
    # ---------------- تب ورودی چت بات ----------------
    chatbot_frame = ttk.Frame(notebook, padding=10)
    notebook.add(chatbot_frame, text="ورودی چت بات")
    
    chatbot_prompt_label = ttk.Label(chatbot_frame, text="ورودی پرامت چت بات (برای تولید متن پاورپوینت):")
    chatbot_prompt_label.grid(row=0, column=0, sticky="w", pady=(0,5))
    
    chatbot_prompt_input = ScrolledText(chatbot_frame, height=5, width=100, font=("Segoe UI", 11), wrap="word")
    chatbot_prompt_input.grid(row=1, column=0, padx=10, pady=10, sticky="nsew")
    add_context_menu(chatbot_prompt_input)
    
    chatbot_output_label = ttk.Label(chatbot_frame, text="متن تولید شده توسط چت بات:")
    chatbot_output_label.grid(row=2, column=0, sticky="w", pady=(10,5))
    
    chatbot_output = ScrolledText(chatbot_frame, height=10, width=100, font=("Segoe UI", 11), wrap="word")
    chatbot_output.grid(row=3, column=0, padx=10, pady=10, sticky="nsew")
    add_context_menu(chatbot_output)
    
    chatbot_frame.grid_rowconfigure(1, weight=0)
    chatbot_frame.grid_rowconfigure(3, weight=1)
    chatbot_frame.grid_columnconfigure(0, weight=1)
    
    chatbot_buttons_frame = ttk.Frame(chatbot_frame, padding=10)
    chatbot_buttons_frame.grid(row=4, column=0, sticky="ew")
    
    def generate_from_chatbot():
        root.config(cursor="wait")
        root.update()
        if not check_app_state():
            messagebox.showerror("غیر فعال", "برنامه فعلا در دسترس نیست.\nلطفاً بعداً امتحان کنید.")
            root.config(cursor="")
            return
        prompt = chatbot_prompt_input.get("1.0", "end").strip()
        if not prompt:
            messagebox.showerror("خطا", "پرامت چت بات خالی است!")
            root.config(cursor="")
            return
        try:
            client = genai.Client(api_key="AIzaSyA1hN1uGimOy-IxttnWB9WSvDXx_uvlBok")
            response = client.models.generate_content(
                model="gemini-2.0-flash",
                contents=prompt,
            )
            generated_text = response.text
            chatbot_output.delete("1.0", "end")
            chatbot_output.insert("1.0", generated_text)
        except Exception as e:
            messagebox.showerror("خطا", f"در تولید متن از چت بات مشکلی به وجود آمد:\n{e}")
        root.config(cursor="")
    
    generate_chatbot_button = ttk.Button(chatbot_buttons_frame, text="تولید متن از چت بات", command=generate_from_chatbot)
    generate_chatbot_button.pack(side='left', padx=5)
    CreateToolTip(generate_chatbot_button, "متن را از چت بات دریافت می‌کند")
    
    def clear_chatbot_text():
        chatbot_prompt_input.delete("1.0", "end")
        chatbot_output.delete("1.0", "end")
    
    clear_chatbot_button = ttk.Button(chatbot_buttons_frame, text="پاکسازی ورودی و خروجی چت بات", command=clear_chatbot_text)
    clear_chatbot_button.pack(side='left', padx=5)
    CreateToolTip(clear_chatbot_button, "متن‌های چت بات را پاک می‌کند")
    
    def transfer_text_from_chatbot_to_manual():
        text = chatbot_output.get("1.0", "end").strip()
        if text:
            text_input.delete("1.0", "end")
            text_input.insert("1.0", text)
            messagebox.showinfo("انتقال متن", "متن از تب چت بات به تب ورودی منتقل شد.")
        else:
            messagebox.showerror("خطا", "متن تولید شده چت بات خالی است!")
    
    transfer_from_chatbot_button = ttk.Button(chatbot_buttons_frame, text="انتقال متن تولیدی از چت بات به تب ورودی", command=transfer_text_from_chatbot_to_manual)
    transfer_from_chatbot_button.pack(side='left', padx=5)
    CreateToolTip(transfer_from_chatbot_button, "متن تولید شده را به تب ورودی منتقل می‌کند")
    
    # ---------------- تب راهنما و درباره ----------------
    guide_frame = ttk.Frame(notebook, padding=10)
    notebook.add(guide_frame, text="راهنما")
    
    guide_notice_label = ttk.Label(guide_frame, text="برای مشاهده راهنمای استفاده، این تب را کلیک کنید")
    guide_notice_label.pack(pady=20)
    
    guide_text = (
        "راهنمای استفاده از پاورپوینت ساز هوشمند :👇 \n\n"
        "این پرامتی که در تب ورودی متن میبینی ، پرامت ساخت متن پاورپوینت  با فرمت ورودی مورد انتظاره و باید موضوع و تعداد اسلاید را وارد کنی و سپس با کلیک روی دکمه (انتقال متن به ورودی چت بات)، پرامت اماده شده با موضوع و تعداد اسلاید جایگذاری شده بصورت خودکار، در ورودی چت بات قرار بگیره و با زدن روی دکمه (تولید متن از چت بات)، متن نوشتاری پاورپویت ساخته بشه و با زدن دکمه(انتقال متن به تب ورودی) ، متن نوشتاری نهایی برای ساخت پاور پوینت ار آن آماده بشود. : \n\n"
        "اولین خط پرامت  که شامل دو [] هستش را دستکاری نکن چون برای برنامه بهینه شده\n\n\n\n\n"
        "سایر توضیحات : \n\n\n"
        "1. ورودی متن: در این تب، شما می‌توانید متن کامل پاورپوینت خود را به صورت دستی هم وارد کنی. متن باید دارای فرمت مشخص شده باشد.\n\n"
        "2. ورودی چت بات: در این تب، می‌توانید از هوش مصنوعی برای تولید متن پاورپوینت استفاده کنید.\n\n"
        "3. درباره: در این تب، اطلاعات کلی در مورد نرم‌افزار و ویژگی‌های آن ارائه می‌شود.\n\n"
        "4. تولید و آپلود: پس از وارد کردن متن در تب 'ورودی متن' یا انتقال متن از تب 'ورودی چت بات'، بر روی دکمه 'تولید و آپلود پاورپوینت' کلیک کنید تا اسلایدها ساخته شده و فایل به صورت آنلاین آپلود شود. لینک دانلود در همان تب نمایش داده می‌شود.\n\n"
        "در صورت بروز مشکل یا نیاز به راهنمایی بیشتر، لطفاً با تیم پشتیبانی تماس بگیرید."
    )
    
    about_frame = ttk.Frame(notebook, padding=10)
    notebook.add(about_frame, text="درباره")
    
    about_text = (
        "پاورپوینت ساز هوشمند\n"
        "----------------------------------------------------\n"
        "این نرم‌افزار پیشرفته یک سیستم یکپارچه و هوشمند برای تولید و آپلود ارائه‌های پاورپوینت است. با بکارگیری فناوری‌های نوین "
        "این ابزار توانایی تبدیل متون ورودی (دستی یا از طریق چت‌بات) "
        "به اسلایدهای حرفه‌ای و باکیفیت را داراست.\n\n"
        "ویژگی‌های کلیدی:\n"
        "• تولید خودکار اسلایدهای استاندارد\n"
        "• آپلود آنلاین و ارائه لینک دانلود\n"
        "• واسط کاربری کاربرپسند و امکانات ویرایش پیشرفته\n"
        "• پشتیبانی دقیق از متون فارسی\n"
        "• فعال یا غیر فعالسازی از راه دور عملکرد کلی برنامه، توسط سازنده\n\n"
        "این ابزار برای ارائه‌های تجاری، آموزشی، پژوهشی و سایر کاربردهای حرفه‌ای طراحی شده است. برای دریافت پشتیبانی و ارائه نظرات، لطفاً با تیم ما تماس بگیرید."
        "\nAli Rashidi ❤️"
        "\nt.me/WriteYourWay"
    )
    
    about_prompt = (
        "این پرامتی که در تب ورودی متن میبینی ، پرامت ساخت متن پاورپوینت  با فرمت ورودی مورد انتظاره و باید موضوع و تعداد اسلاید را وارد کنی و سپس با کلیک روی دکمه (انتقال متن به ورودی چت بات)، پرامت اماده شده با موضوع و تعداد اسلاید جایگذاری شده بصورت خودکار، در ورودی چت بات قرار بگیره و با زدن روی دکمه (تولید متن از چت بات)، متن نوشتاری پاورپویت ساخته بشه و با زدن دکمه(انتقال متن به تب ورودی) ، متن نوشتاری نهایی برای ساخت پاور پوینت ار آن آماده بشود. : \n\n"
        "اولین خط پرامت  که شامل دو [] هستش را دستکاری نکن چون برای برنامه بهینه شده\n\n"
        "حتما از فیلترشکن یا تحریم شکن قوی و پایدار برای کارکرد صحیح چت بات، استفاده کن\n\n"
        " حتما برو به *تب راهنما* تا راهنما بیشتر را بخونی و بدونی چه کاری باید انجام باید بدی *\n"
        "-------------------------\n\n"
        "👈👈👈پرامت👇👇👇 \n\n\n\n\n"
    )
    
    def on_tab_changed(event):
        global text_tab_shown
        selected = event.widget.tab(event.widget.select(), "text")
        if selected == "راهنما":
            messagebox.showinfo("راهنما", guide_text)
            event.widget.select(manual_frame)
        elif selected == "درباره":
            messagebox.showinfo("درباره", about_text)
            event.widget.select(manual_frame)
        elif selected == "ورودی متن" and not text_tab_shown:
            # استفاده از after برای تأخیر در نمایش اعلان
            root.after(100, lambda: show_initial_notification())

    def show_initial_notification():
        global text_tab_shown
        messagebox.showinfo("توجه", about_prompt)
        text_tab_shown = True
        # برگرداندن تمرکز به فیلد موضوع
        subject_entry.focus_set()
    
    notebook.bind("<<NotebookTabChanged>>", on_tab_changed)
    
    # ---------------- نوار وضعیت ----------------
    status_var = tk.StringVar()
    status_var.set("آماده")
    status_bar = ttk.Label(root, textvariable=status_var, relief="sunken", anchor="w", font=("Segoe UI", 10),
                           background="#283655", foreground="white")
    status_bar.pack(side="bottom", fill="x")
    
    root.mainloop()

if __name__ == "__main__":
    main()
