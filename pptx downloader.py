import io
import os
import time
import requests
from pptx import Presentation

def create_pptx_in_memory():
    """
    Create a PowerPoint file with 10 slides based on predefined content
    and return it as an in-memory BytesIO object.
    """
    prs = Presentation()
    slides_data = [
        {
            "title": "Introducing PromptX Startup: Vision & Mission",
            "content": "- Brief introduction to the startup\n- Long-term vision and key objectives\n- Unique value proposition for entrepreneurs"
        },
        {
            "title": "Identifying the Problem: Market Needs & Opportunities",
            "content": "- Description of current issues in idea generation and documentation\n- Analysis of customer needs and challenges\n- Opportunities in the startup market"
        },
        {
            "title": "Our Solution: PromptX Platform",
            "content": "- Explanation of how PromptX solves the problems\n- Key features (idea generation, pitch deck, and technical document)\n- Competitive advantages over traditional methods"
        },
        {
            "title": "Business Model & Revenue Strategies",
            "content": "- Explanation of revenue models (monthly subscription, pay-per-use, consultancy services)\n- Pricing strategy and profitability\n- Plans for market expansion"
        },
        {
            "title": "Target Market & Competitive Landscape",
            "content": "- Analysis of the target market and customer segmentation\n- Competitor analysis, strengths and weaknesses\n- PromptX's competitive edge in the market"
        },
        {
            "title": "Marketing & Sales Strategies",
            "content": "- Distribution and advertising channels\n- Customer acquisition plans and digital marketing strategies\n- Content strategy and brand awareness initiatives"
        },
        {
            "title": "Our Team: Experts Behind PromptX",
            "content": "- Introduction of key team members and their expertise\n- Organizational structure and key roles\n- Team strengths in project execution and innovation"
        },
        {
            "title": "Financial Outlook: Growth & Profitability",
            "content": "- Financial charts including revenue, costs, and profitability forecasts\n- Key financial milestones and growth plans\n- Investment analysis and ROI"
        },
        {
            "title": "Roadmap & Development Plan",
            "content": "- Timeline for developing the MVP and future steps\n- Short-term and long-term plans\n- Key objectives for different development phases"
        },
        {
            "title": "Investment Opportunity: Join PromptX",
            "content": "- Summary of the investment opportunities in the startup\n- Financial needs and growth outlook\n- Invitation to collaborate and participate in the success of PromptX"
        }
    ]

    for slide in slides_data:
        slide_layout = prs.slide_layouts[1]  # Use the "Title and Content" layout
        current_slide = prs.slides.add_slide(slide_layout)
        current_slide.shapes.title.text = slide["title"]
        current_slide.placeholders[1].text = slide["content"]

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

def upload_file(file_bytes, file_name, upload_url):
    """
    Upload the in-memory file via the proxy using multipart/form-data.
    """
    files = {
        "file": (file_name, file_bytes, "application/vnd.openxmlformats-officedocument.presentationml.presentation")
    }
    data = {
        "destination_folder": "store_html_report"
    }
    headers = {
        "User-Agent": "Mozilla/5.0",
        "Expect": ""  # Remove the default "Expect" header
    }
    
    # Send POST request without auto-following redirects
    response = requests.post(upload_url, files=files, data=data, headers=headers, allow_redirects=False)
    
    # If server returns a redirect, follow it manually
    if response.status_code in (301, 302, 303, 307, 308):
        redirect_url = response.headers.get("Location")
        if redirect_url:
            print("Redirected to:", redirect_url)
            response = requests.post(redirect_url, files=files, data=data, headers=headers)
    
    if response.status_code == 200:
        try:
            resp_data = response.json()
            download_link = resp_data.get("download_link", "").strip()
            if download_link:
                print("Download link received:", download_link)
                return download_link
            else:
                download_link = response.text.strip()
                print("Download link (from response text):", download_link)
                return download_link
        except Exception as e:
            print("Error processing JSON response:", e)
            download_link = response.text.strip()
            return download_link
    else:
        print("File upload failed. Status code:", response.status_code)
        print("Response text:", response.text)
        return None

def save_download_link(download_link, txt_file):
    """
    Save the final download link into a text file.
    """
    with open(txt_file, 'w', encoding='utf-8') as f:
        f.write(download_link)
    print(f"Download link saved in file '{txt_file}'.")

def main():
    # Generate a unique filename for the PowerPoint file
    timestamp = int(time.time())
    pptx_filename = f"promptx_pitch_deck_{timestamp}.pptx"
    
    # Create the PowerPoint file in memory (without saving to disk)
    pptx_bytes = create_pptx_in_memory()
    
    # Proxy URL for file upload
    upload_url = "http://www.rash32.ir/python/micropython/file_uploader_proxy.php"
    
    # Upload the file and get the download link
    download_link = upload_file(pptx_bytes, pptx_filename, upload_url)
    
    # Save the download link to a text file if successful
    if download_link:
        save_download_link(download_link, "download_link.txt")
    else:
        print("There was an issue retrieving the download link.")

if __name__ == "__main__":
    main()
