import io
import os
import re
import time
import requests
import tkinter as tk
from tkinter import ttk, messagebox
from pptx import Presentation
from pptx.enum.text import PP_ALIGN  # جهت تنظیم تراز متن
from pptx.oxml.xmlchemy import OxmlElement  # جهت دستکاری XML برای حذف bullets
import webbrowser

# کتابخانه‌های مورد نیاز جهت پردازش صحیح متون فارسی
import arabic_reshaper
from bidi.algorithm import get_display

def is_rtl(text):
    """
    بررسی می‌کند که آیا متن شامل کاراکترهای فارسی/عربی است یا خیر.
    """
    return bool(re.search(r'[\u0600-\u06FF]', text))

def process_text(text):
    """
    اگر متن شامل کاراکترهای راست‌چین (فارسی/عربی) باشد،
    ابتدا با arabic_reshaper شکل‌دهی و سپس با get_display پردازش می‌شود.
    در غیر این صورت، همان متن اصلی برگردانده می‌شود.
    """
    if is_rtl(text):
        reshaped_text = arabic_reshaper.reshape(text)
        bidi_text = get_display(reshaped_text)
        return bidi_text
    return text

def remove_bullets(text_frame):
    """
    از درج bullet (علامت گلوله‌ای) در هر پاراگراف جلوگیری می‌کند.
    این تابع به سطح XML دسترسی پیدا کرده و عنصر <a:buNone/> را اضافه می‌کند.
    """
    for paragraph in text_frame.paragraphs:
        pPr = paragraph._p.get_or_add_pPr()
        buNone = OxmlElement('a:buNone')
        pPr.append(buNone)

def parse_slides(text):
    """
    متن ورودی شامل اسلایدهای مختلف (با فرمت مشخص‌شده) را دریافت کرده و
    لیستی از دیکشنری‌های حاوی "title" (عنوان) و "content" (محتوا) را برمی‌گرداند.
    
    فرمت مورد انتظار:
    
    اسلاید ۱: مقدمه و چشم‌انداز  
    عنوان اسلاید: معرفی استارتاپ PromptX: چشم‌انداز و ماموریت  
    محتوا:  
    معرفی مختصر استارتاپ  
    چشم‌انداز بلندمدت و اهداف کلیدی  
    ارزش پیشنهادی منحصر به فرد برای کارآفرینان  
    
    توجه: از درج bullet یا جداکننده‌های اضافی استفاده نکنید.
    """
    blocks = re.split(r"(?=اسلاید\s*\d+:)", text)
    slides = []
    for block in blocks:
        block = block.strip()
        if not block:
            continue
        m_title = re.search(r"عنوان\s+اسلاید\s*:\s*(.+)", block)
        m_content = re.search(r"محتوا\s*:\s*(.+)", block, re.DOTALL)
        title = m_title.group(1).strip() if m_title else ""
        content = m_content.group(1).strip() if m_content else ""
        slides.append({"title": title, "content": content})
    return slides

def create_pptx_in_memory(slides_data):
    """
    بر اساس لیست اسلایدهای دریافتی، یک فایل پاورپوینت به صورت in‑memory تولید می‌کند 
    و شیء BytesIO را برمی‌گرداند. در این تابع:
      - متون عنوان و محتوا پردازش و جهت‌چین می‌شوند،
      - از درج bullet جلوگیری می‌شود.
    """
    prs = Presentation()
    for slide in slides_data:
        slide_layout = prs.slide_layouts[1]  # قالب "عنوان و محتوا"
        current_slide = prs.slides.add_slide(slide_layout)
        
        # پردازش و تنظیم عنوان
        title_shape = current_slide.shapes.title
        processed_title = process_text(slide["title"])
        title_shape.text = processed_title
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.RIGHT if is_rtl(slide["title"]) else PP_ALIGN.LEFT
        remove_bullets(title_shape.text_frame)
        
        # پردازش و تنظیم محتوا
        content_shape = current_slide.placeholders[1]
        processed_content = process_text(slide["content"])
        content_shape.text = processed_content
        for paragraph in content_shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.RIGHT if is_rtl(slide["content"]) else PP_ALIGN.LEFT
        remove_bullets(content_shape.text_frame)
    
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

def upload_file(file_bytes, file_name, upload_url):
    """
    فایل تولید شده به صورت in‑memory را از طریق پروکسی (با استفاده از multipart/form-data) آپلود می‌کند.
    از محتوای فایل در یک متغیر استفاده می‌شود تا در صورت redirect، شیء BytesIO جدید ساخته شود.
    در صورت موفقیت، لینک دانلود نهایی را برمی‌گرداند.
    """
    file_content = file_bytes.getvalue()
    
    def post_request(url):
        file_obj = io.BytesIO(file_content)
        files = {
            "file": (file_name, file_obj, "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        }
        data = {
            "destination_folder": "store_html_report"
        }
        headers = {
            "User-Agent": "Mozilla/5.0",
            "Expect": ""
        }
        return requests.post(url, files=files, data=data, headers=headers, allow_redirects=False)
    
    response = post_request(upload_url)
    
    if response.status_code in (301, 302, 303, 307, 308):
        redirect_url = response.headers.get("Location")
        if redirect_url:
            print("Redirected to:", redirect_url)
            response = post_request(redirect_url)
    
    if response.status_code == 200:
        try:
            resp_data = response.json()
            download_link = resp_data.get("download_link", "").strip()
            print("Download link received:", download_link)
            return download_link
        except Exception as e:
            print("Error processing JSON response:", e)
            return response.text.strip()
    else:
        print("File upload failed. Status code:", response.status_code)
        print("Response text:", response.text)
        return None

def save_download_link(download_link, txt_file):
    """
    لینک دانلود نهایی را در یک فایل متنی ذخیره می‌کند.
    """
    with open(txt_file, 'w', encoding='utf-8') as f:
        f.write(download_link)
    print(f"Download link saved in file '{txt_file}'.")

def main():
    root = tk.Tk()
    root.title("تولید و آپلود پاورپوینت - برنامه علی رشیدی")
    root.geometry("1000x750")
    root.configure(bg="#f5f5f5")
    
    style = ttk.Style()
    style.theme_use('clam')
    style.configure("TLabel", background="#f5f5f5", font=("Tahoma", 11))
    style.configure("TButton", font=("Tahoma", 11))
    
    notebook = ttk.Notebook(root)
    notebook.pack(fill='both', expand=True)
    
    # تب "ورودی"
    input_frame = ttk.Frame(notebook, padding=10)
    notebook.add(input_frame, text="ورودی")
    
    input_label = ttk.Label(input_frame, text="متن کامل پاورپوینت (با فرمت مورد انتظار) را وارد کنید:", font=("Tahoma", 12))
    input_label.pack(pady=10, anchor="w")
    
    text_input = tk.Text(input_frame, height=25, width=100, font=("Tahoma", 11))
    text_input.pack(padx=10, pady=10)
    
    # نمونه متن همراه با راهنمای فرمت که در کادر ورودی قرار گرفته است.
    sample_text = (
        "فرمت ورودی مورد انتظار:\n"
        "-------------------------\n"
        "• هر اسلاید با عبارت \"اسلاید شماره:\" آغاز می‌شود. به عنوان مثال:\n"
        "   اسلاید ۱: مقدمه و چشم‌انداز\n\n"
        "• سپس خطی با عبارت \"عنوان اسلاید:\" نوشته می‌شود که عنوان اسلاید را مشخص می‌کند. به عنوان مثال:\n"
        "   عنوان اسلاید: معرفی استارتاپ PromptX: چشم‌انداز و ماموریت\n\n"
        "• در ادامه، بخش \"محتوا:\" قرار می‌گیرد که شامل متن اصلی اسلاید است. به عنوان مثال:\n"
        "   محتوا:\n"
        "   معرفی مختصر استارتاپ\n"
        "   چشم‌انداز بلندمدت و اهداف کلیدی\n"
        "   ارزش پیشنهادی منحصر به فرد برای کارآفرینان\n\n"
        "توجه: از درج bullet یا جداکننده‌های اضافی استفاده نکنید.\n\n"
        "-------------------------\n\n"
        "مثال:\n\n"
        "اسلاید ۱: مقدمه و چشم‌انداز\n"
        "عنوان اسلاید: معرفی استارتاپ PromptX: چشم‌انداز و ماموریت\n"
        "محتوا:\n"
        "معرفی مختصر استارتاپ\n"
        "چشم‌انداز بلندمدت و اهداف کلیدی\n"
        "ارزش پیشنهادی منحصر به فرد برای کارآفرینان\n\n"
        "اسلاید ۲: مشکل و نیاز بازار\n"
        "عنوان اسلاید: شناسایی مشکل: نیازهای بازار و فرصت‌های موجود\n"
        "محتوا:\n"
        "توضیح مشکلات موجود در حوزه تولید ایده و مستندسازی\n"
        "تحلیل دقیق نیازهای مشتریان و چالش‌های فعلی\n"
        "فرصت‌های موجود در بازار استارتاپ‌های نوپا\n"
    )
    text_input.insert("1.0", sample_text)
    
    # دکمه "پاکسازی" جهت پاک کردن محتویات کادر ورودی
    def clear_input():
        text_input.delete("1.0", "end")
    
    clear_button = ttk.Button(input_frame, text="پاکسازی متن", command=clear_input)
    clear_button.pack(pady=5, anchor="e", padx=10)
    
    # منوی راست کلیکی جهت امکان Paste
    popup_menu = tk.Menu(text_input, tearoff=0)
    popup_menu.add_command(label="Paste", command=lambda: text_input.event_generate("<<Paste>>"))
    
    def show_popup(event):
        popup_menu.tk_popup(event.x_root, event.y_root)
    
    text_input.bind("<Button-3>", show_popup)
    text_input.bind("<Control-v>", lambda e: text_input.event_generate("<<Paste>>"))
    text_input.bind("<Control-V>", lambda e: text_input.event_generate("<<Paste>>"))
    text_input.bind("<Control-c>", lambda e: text_input.event_generate("<<Copy>>"))
    text_input.bind("<Control-C>", lambda e: text_input.event_generate("<<Copy>>"))
    text_input.bind("<Control-x>", lambda e: text_input.event_generate("<<Cut>>"))
    text_input.bind("<Control-X>", lambda e: text_input.event_generate("<<Cut>>"))
    
    # تب "درباره"
    about_frame = ttk.Frame(notebook, padding=10)
    notebook.add(about_frame, text="درباره")
    about_text = (
        "برنامه علی رشیدی\n"
        "-------------------------\n"
        "این برنامه ابزاری قدرتمند برای تولید و آپلود پاورپوینت به صورت تعاملی و چندزبانه است.\n"
        "با بهره‌گیری از فناوری‌های نوین مانند python-pptx، arabic_reshaper و python-bidi،\n"
        "امکان تولید اسلایدهای حرفه‌ای، بی‌نقص و مطابق با استانداردهای نگارشی فارسی و انگلیسی فراهم شده است.\n\n"
        "این برنامه توسط علی رشیدی، یک کارآفرین و متخصص با بیش از 2 سال تجربه در حوزه فناوری و استارتاپ‌ها\n"
        "طراحی و توسعه یافته است."
    )
    about_label = ttk.Label(about_frame, text=about_text, font=("Tahoma", 11), justify="left")
    about_label.pack(padx=10, pady=10, anchor="w")
    
    # تب "تنظیمات"
    settings_frame = ttk.Frame(notebook, padding=10)
    notebook.add(settings_frame, text="تنظیمات")
    settings_text = (
        "تنظیمات فعلی:\n"
        "-------------------------\n"
        "در حال حاضر تنظیمات پیش‌فرض استفاده می‌شود.\n"
        "در نسخه‌های آینده، امکانات بیشتری مانند تغییر فونت، رنگ‌ها و سایر گزینه‌ها اضافه خواهد شد."
    )
    settings_label = ttk.Label(settings_frame, text=settings_text, font=("Tahoma", 11), justify="left")
    settings_label.pack(padx=10, pady=10, anchor="w")
    
    # فریم عملیات در تب "ورودی"
    input_action_frame = ttk.Frame(input_frame, padding=10)
    input_action_frame.pack(fill='x')
    
    download_link_var = tk.StringVar()
    
    def process_input():
        raw_text = text_input.get("1.0", "end").strip()
        if not raw_text:
            messagebox.showerror("خطا", "متن ورودی خالی است!")
            return
        slides_data = parse_slides(raw_text)
        if not slides_data:
            messagebox.showerror("خطا", "نتوانستیم اسلایدها را از متن استخراج کنیم.\nلطفاً فرمت ورودی را بررسی کنید.")
            return
        pptx_bytes = create_pptx_in_memory(slides_data)
        timestamp = int(time.time())
        file_name = f"promptx_pitch_deck_{timestamp}.pptx"
        upload_url = "http://www.rash32.ir/python/micropython/file_uploader_proxy.php"
        download_link = upload_file(pptx_bytes, file_name, upload_url)
        if download_link:
            download_link_var.set(download_link)
            messagebox.showinfo("موفقیت", f"پاورپوینت با موفقیت تولید و آپلود شد!\nلینک دانلود:\n{download_link}")
        else:
            messagebox.showerror("خطا", "آپلود فایل با مشکل مواجه شد.")
    
    generate_button = ttk.Button(input_action_frame, text="تولید و آپلود پاورپوینت", command=process_input)
    generate_button.pack(side='left', padx=5)
    
    download_link_label = ttk.Label(input_action_frame, textvariable=download_link_var, foreground="blue", cursor="hand2", font=("Tahoma", 11, "underline"))
    download_link_label.pack(side='left', padx=10)
    download_link_label.bind("<Button-1>", lambda e: webbrowser.open(download_link_var.get()))
    
    root.mainloop()

if __name__ == "__main__":
    main()
